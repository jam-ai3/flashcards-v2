from pptx import Presentation
import os
from io import BytesIO


def get_text_from_pptx(input_data):
    """
    Extract text from a PPTX file
    Args:
        input_data: Either a file path (str) or binary data (bytes)
    Returns:
        str: Extracted text from the presentation
    """
    try:
        presentation = Presentation(BytesIO(input_data))
        all_text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        all_text += text + "\n"  # Add newline between text blocks

        return all_text.strip()  # Remove trailing newline
    except Exception as e:
        print(f"Error in get_text_from_pptx: {str(e)}")
        raise


if __name__ == "__main__":
    # Get the directory where this script is located
    curr_path = os.path.dirname(os.path.abspath(__file__))
    pptx_name = "pptx4.pptx"
    file_path = os.path.join(curr_path, pptx_name)
    text = get_text_from_pptx(file_path)
    print(text)
